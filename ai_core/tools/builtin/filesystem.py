"""Filesystem tools — safe read/write for all file types in the workspace.

read_file auto-detects format: text, PDF, Word, Excel, PPT, and images (returns base64).
write_file auto-detects format: text for code/docs, base64 decode for binary files.
All paths are sandboxed to the workspace root.
"""

import base64
import io
import os
import re
import shutil
import fnmatch
from tools.registry import get_tool_registry
from tools.schema import PermissionLevel, ToolDef


_IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "bmp", "webp"}
_TEXT_EXTS = {"txt", "py", "js", "ts", "jsx", "tsx", "vue", "html", "css", "json",
              "yaml", "yml", "toml", "md", "rst", "csv", "xml", "svg", "sh", "bat",
              "ps1", "ini", "cfg", "conf", "log", "env", "gitignore", "dockerfile",
              "c", "cpp", "h", "hpp", "rs", "go", "java", "kt", "swift", "rb", "php",
              "sql", "r", "m", "mm", "proto", "gradle", "lock"}
_BINARY_EXTS = {"png", "jpg", "jpeg", "gif", "bmp", "webp", "pdf", "docx",
                "xlsx", "xls", "pptx", "zip", "tar", "gz", "exe", "dll",
                "so", "dylib", "bin", "dat", "db", "sqlite"}


def _workspace_root() -> str:
    import os as _os
    default = _os.path.join(_os.path.dirname(_os.path.dirname(_os.path.dirname(_os.path.abspath(__file__)))), "game_workspace")
    return _os.environ.get("TOOL_WORKSPACE", default)


def _resolve_path(filename: str, scope: str = "workspace") -> str:
    """Resolve a file path.

    ``workspace`` (default) confines the path to the workspace sandbox and
    rejects traversal. ``system`` allows any path on disk — this is gated by
    the ``file.*.system`` capability at the execution layer, so the sandbox is
    only lifted when the operator has explicitly authorized system file access.
    """
    if scope == "system":
        return os.path.realpath(filename)
    root = os.path.realpath(_workspace_root())
    candidate = os.path.realpath(os.path.join(root, filename))
    common = os.path.commonpath([root, candidate])
    if common != root:
        raise ValueError(f"Path traversal denied: {filename!r}")
    return candidate


def _resolve_base(path: str, scope: str) -> str:
    """Resolve a directory base for list/search operations."""
    if scope == "system":
        return os.path.realpath(path) if path else os.getcwd()
    root = os.path.realpath(_workspace_root())
    base = os.path.realpath(os.path.join(root, path)) if path else root
    if os.path.commonpath([root, base]) != root:
        raise ValueError(f"Path traversal denied: {path!r}")
    return base


def _display_root(scope: str) -> str:
    """The reference path used to render relative paths in results."""
    return os.getcwd() if scope == "system" else os.path.realpath(_workspace_root())


def _ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def _read_text(path: str, offset: int, limit: int) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        lines = f.readlines()
    if offset:
        lines = lines[offset:]
    if limit and len(lines) > limit:
        lines = lines[:limit]
        lines.append(f"\n... (truncated at {limit} lines)")
    return "".join(lines)


def _read_image(path: str, filename: str) -> str:
    try:
        from PIL import Image
        img = Image.open(path)
        w, h = img.size
        fmt = img.format or _ext(filename).upper()
        with open(path, "rb") as f:
            raw = f.read()
        b64 = base64.b64encode(raw).decode()
        mime_map = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                    "gif": "image/gif", "bmp": "image/bmp", "webp": "image/webp"}
        mime = mime_map.get(_ext(filename), "application/octet-stream")
        return f"[image,file={filename}] {filename} ({w}x{h}, {len(raw)} bytes)\n[image,base64={b64}]"
    except ImportError:
        return f"Error: PIL/Pillow not installed. Cannot read image: {filename}"
    except Exception as e:
        return f"Error reading image {filename}: {e}"


def _read_pdf(path: str, filename: str) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
        if not pages:
            return f"{filename}: PDF appears to have no extractable text ({len(reader.pages)} pages)."
        return "\n\n".join(pages)
    except ImportError:
        return "Error: PyPDF2 not installed."
    except Exception as e:
        return f"Error reading PDF {filename}: {e}"


def _read_docx(path: str, filename: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        if not paras:
            return f"{filename}: Word document appears empty."
        return "\n".join(paras)
    except ImportError:
        return "Error: python-docx not installed."
    except Exception as e:
        return f"Error reading Word document {filename}: {e}"


def _read_xlsx(path: str, filename: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        out = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                out.append(f"--- Sheet: {name} ---")
                out.extend(rows)
        wb.close()
        return "\n".join(out) if out else f"{filename}: Excel file appears empty."
    except ImportError:
        return "Error: openpyxl not installed."
    except Exception as e:
        return f"Error reading Excel file {filename}: {e}"


def _read_xls(path: str, filename: str) -> str:
    try:
        import xlrd
        wb = xlrd.open_workbook(path)
        out = []
        for name in wb.sheet_names():
            ws = wb.sheet_by_name(name)
            rows = []
            for r in range(ws.nrows):
                cells = [str(ws.cell_value(r, c)) if ws.cell_value(r, c) != "" else "" for c in range(ws.ncols)]
                if any(c for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                out.append(f"--- Sheet: {name} ---")
                out.extend(rows)
        return "\n".join(out) if out else f"{filename}: Excel file appears empty."
    except ImportError:
        return "Error: xlrd not installed."
    except Exception as e:
        return f"Error reading Excel file {filename}: {e}"


def _read_pptx(path: str, filename: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
        return "\n\n".join(slides) if slides else f"{filename}: PPT has no text content."
    except ImportError:
        return "Error: python-pptx not installed."
    except Exception as e:
        return f"Error reading PPT {filename}: {e}"


async def _read_file(filename: str, offset: int = 0, limit: int = 2000, scope: str = "workspace") -> str:
    path = _resolve_path(filename, scope)
    ext = _ext(filename)

    try:
        if ext in _IMAGE_EXTS:
            return _read_image(path, filename)
        if ext == "pdf":
            return _read_pdf(path, filename)
        if ext == "docx":
            return _read_docx(path, filename)
        if ext == "xlsx":
            return _read_xlsx(path, filename)
        if ext == "xls":
            return _read_xls(path, filename)
        if ext == "pptx":
            return _read_pptx(path, filename)
        return _read_text(path, offset, limit)
    except FileNotFoundError:
        return f"Error: file not found — {filename}"
    except IsADirectoryError:
        return f"Error: {filename} is a directory"
    except ValueError as e:
        return str(e)
    except OSError as e:
        return f"Error reading {filename}: {e}"


async def _list_directory(path: str = "", scope: str = "workspace") -> str:
    root = _workspace_root()
    if scope != "system":
        os.makedirs(root, exist_ok=True)
    try:
        target = _resolve_base(path, scope)
    except ValueError as e:
        return f"Error: {e}"
    try:
        names = sorted(os.listdir(target))
    except NotADirectoryError:
        return f"Error: {path} is not a directory"
    except FileNotFoundError:
        return f"Error: directory not found — {path}"
    except OSError as e:
        return f"Error listing directory: {e}"
    lines = [f"Contents of {path or '(root)'}:"]
    for n in names:
        full = os.path.join(target, n)
        if os.path.isdir(full):
            lines.append(f"  [DIR]  {n}/")
        else:
            size = os.path.getsize(full)
            lines.append(f"  [FILE] {n}  ({size} bytes)")
    if len(lines) == 1:
        lines.append("  (empty)")
    return "\n".join(lines)


async def _search_files(pattern: str, path: str = "", scope: str = "workspace") -> str:
    try:
        base = _resolve_base(path, scope)
    except ValueError as e:
        return f"Error: {e}"
    ref = _display_root(scope)
    matches = []
    try:
        for dirpath, _dirnames, filenames in os.walk(base):
            for fn in filenames:
                if fnmatch.fnmatch(fn, pattern):
                    rel = os.path.relpath(os.path.join(dirpath, fn), ref)
                    size = os.path.getsize(os.path.join(dirpath, fn))
                    matches.append(f"  {rel}  ({size} bytes)")
    except OSError as e:
        return f"Error searching files: {e}"
    if not matches:
        return f"No files matching {pattern!r} found."
    matches.sort()
    return f"Found {len(matches)} file(s) matching {pattern!r}:\n" + "\n".join(matches)


async def _search_content(pattern: str, glob: str = "*", path: str = "", max_results: int = 30, scope: str = "workspace") -> str:
    try:
        base = _resolve_base(path, scope)
    except ValueError as e:
        return f"Error: {e}"
    ref = _display_root(scope)
    results = []
    try:
        compiled = re.compile(pattern, re.IGNORECASE)
    except re.error as e:
        return f"Error: invalid regex pattern — {e}"
    try:
        for dirpath, _dirnames, filenames in os.walk(base):
            for fn in filenames:
                if not fnmatch.fnmatch(fn, glob):
                    continue
                full = os.path.join(dirpath, fn)
                try:
                    with open(full, "r", encoding="utf-8", errors="replace") as f:
                        for lineno, line in enumerate(f, 1):
                            if compiled.search(line):
                                rel = os.path.relpath(full, ref)
                                results.append(f"  {rel}:{lineno}: {line.strip()[:200]}")
                                if len(results) >= max_results:
                                    break
                except OSError:
                    continue
                if len(results) >= max_results:
                    break
            if len(results) >= max_results:
                break
    except OSError as e:
        return f"Error searching content: {e}"
    if not results:
        return f"No lines matching {pattern!r} found."
    return f"Found {len(results)} match(es) for {pattern!r}:\n" + "\n".join(results)


async def _write_file(filename: str, content: str, scope: str = "workspace") -> str:
    path = _resolve_path(filename, scope)
    ext = _ext(filename)
    os.makedirs(os.path.dirname(path) or _workspace_root(), exist_ok=True)

    if ext in _BINARY_EXTS and ext not in _TEXT_EXTS:
        try:
            raw = base64.b64decode(content)
            with open(path, "wb") as f:
                f.write(raw)
            return f"Wrote {len(raw)} bytes to {filename} (base64 decoded)"
        except Exception as e:
            return f"Error: base64 decode failed for {filename}: {e}. Content must be valid base64 for binary files."
    else:
        try:
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)
            size = os.path.getsize(path)
            return f"Wrote {len(content)} chars ({size} bytes) to {filename}"
        except OSError as e:
            return f"Error writing {filename}: {e}"


async def _edit_file(filename: str, old_string: str, new_string: str, scope: str = "workspace") -> str:
    path = _resolve_path(filename, scope)
    try:
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()
    except FileNotFoundError:
        return f"Error: file not found — {filename}"
    except OSError as e:
        return f"Error reading {filename}: {e}"
    count = content.count(old_string)
    if count == 0:
        return f"Error: old_string not found in {filename}"
    if count > 1:
        return f"Error: old_string found {count} times in {filename}. Provide more context."
    new_content = content.replace(old_string, new_string, 1)
    try:
        with open(path, "w", encoding="utf-8") as f:
            f.write(new_content)
        return f"Replaced 1 occurrence in {filename}. {len(content)} → {len(new_content)} chars."
    except OSError as e:
        return f"Error writing {filename}: {e}"


async def _delete_file(filename: str, scope: str = "workspace") -> str:
    path = _resolve_path(filename, scope)
    try:
        os.remove(path)
        return f"Deleted {filename}"
    except FileNotFoundError:
        return f"Error: file not found — {filename}"
    except IsADirectoryError:
        try:
            shutil.rmtree(path)
            return f"Deleted directory {filename}"
        except OSError as e:
            return f"Error deleting directory {filename}: {e}"
    except OSError as e:
        return f"Error deleting {filename}: {e}"


def register() -> None:
    reg = get_tool_registry()
    reg.register(ToolDef(
        name="read_file",
        description="读取文件内容。自动识别文件类型：纯文本直接返回，PDF/Word/Excel/PPT自动提取文本，图片（png/jpg/gif/bmp/webp）返回base64编码供视觉分析。支持offset/limit分页读取大文件。",
        parameters={
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "工作空间内的文件路径，相对路径。"},
                "offset": {"type": "integer", "description": "从第几行开始读（0开始），仅文本文件有效。"},
                "limit": {"type": "integer", "description": "最多返回多少行（默认2000），仅文本文件有效。"},
                "scope": {"type": "string", "enum": ["workspace", "system"], "description": "路径作用域：workspace=工作空间内（默认，相对路径）；system=工作空间外（需授权，绝对路径）。"},
            },
            "required": ["filename"],
        },
        permission_level=PermissionLevel.READ,
        handler=_read_file,
        group="system",
        category="文件操作",
        guidance="要读某个文件的内容 → read_file",
    ))
    reg.register(ToolDef(
        name="list_directory",
        description="列出工作空间指定目录下的所有文件和子文件夹，显示类型（文件/目录）和大小。",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "子目录路径，相对于工作空间根目录。不填则为根目录。"},
                "scope": {"type": "string", "enum": ["workspace", "system"], "description": "路径作用域：workspace=工作空间内（默认）；system=工作空间外（需授权，path 为绝对路径）。"},
            },
            "required": [],
        },
        permission_level=PermissionLevel.READ,
        handler=_list_directory,
        group="system",
        category="文件操作",
        guidance="想知道工作空间里有什么文件 → list_directory",
    ))
    reg.register(ToolDef(
        name="search_files",
        description="按文件名模式搜索工作空间内的文件，支持通配符（如 *.py、test*.ts、**/*.json）。返回匹配文件的路径和大小。",
        parameters={
            "type": "object",
            "properties": {
                "pattern": {"type": "string", "description": "文件名匹配模式，支持通配符 * 和 **。"},
                "path": {"type": "string", "description": "搜索的子目录，不填则搜索整个工作空间。"},
                "scope": {"type": "string", "enum": ["workspace", "system"], "description": "路径作用域：workspace=工作空间内（默认）；system=工作空间外（需授权，path 为绝对路径）。"},
            },
            "required": ["pattern"],
        },
        permission_level=PermissionLevel.READ,
        handler=_search_files,
        group="system",
        category="文件操作",
        guidance="要找某种类型的文件 → search_files",
    ))
    reg.register(ToolDef(
        name="search_content",
        description="在文件内容中搜索正则表达式匹配，类似grep命令。返回 文件名:行号:匹配行 格式的结果，默认最多返回30条。",
        parameters={
            "type": "object",
            "properties": {
                "pattern": {"type": "string", "description": "正则表达式搜索模式。"},
                "glob": {"type": "string", "description": "文件名过滤，如 *.py。默认 * 匹配所有文件。"},
                "path": {"type": "string", "description": "搜索的子目录。"},
                "scope": {"type": "string", "enum": ["workspace", "system"], "description": "路径作用域：workspace=工作空间内（默认）；system=工作空间外（需授权，path 为绝对路径）。"},
            },
            "required": ["pattern"],
        },
        permission_level=PermissionLevel.READ,
        handler=_search_content,
        group="system",
        category="文件操作",
        guidance="要在文件内容里搜索关键词 → search_content",
    ))
    reg.register(ToolDef(
        name="write_file",
        description="在工作空间写入或覆盖文件。文本文件（.py/.txt/.md等）直接传入字符串内容；二进制文件（.png/.pdf/.docx/.xlsx/.pptx等）需传入base64编码的内容，自动解码后写入。自动创建父目录。",
        parameters={
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "文件路径，相对于工作空间根目录。"},
                "content": {"type": "string", "description": "文件内容。文本文件直接传入文本，二进制文件传入base64编码字符串。"},
                "scope": {"type": "string", "enum": ["workspace", "system"], "description": "路径作用域：workspace=工作空间内（默认）；system=工作空间外（需授权，filename 为绝对路径）。"},
            },
            "required": ["filename", "content"],
        },
        permission_level=PermissionLevel.WRITE,
        handler=_write_file,
        group="system",
        category="文件操作",
        guidance="要创建或覆盖一个文件 → write_file",
    ))
    reg.register(ToolDef(
        name="edit_file",
        description="精确替换文件中的一段文本。old_string必须在文件中唯一出现一次。如果要替换多处或添加新内容，请使用write_file覆盖整个文件。",
        parameters={
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "要编辑的文件路径。"},
                "old_string": {"type": "string", "description": "要被替换的精确文本，必须在文件中唯一。"},
                "new_string": {"type": "string", "description": "替换后的新文本。"},
                "scope": {"type": "string", "enum": ["workspace", "system"], "description": "路径作用域：workspace=工作空间内（默认）；system=工作空间外（需授权，filename 为绝对路径）。"},
            },
            "required": ["filename", "old_string", "new_string"],
        },
        permission_level=PermissionLevel.WRITE,
        handler=_edit_file,
        group="system",
        category="文件操作",
        guidance="要修改文件里的某段内容 → edit_file",
    ))
    reg.register(ToolDef(
        name="delete_file",
        description="删除工作空间中的文件或目录。如果是目录则递归删除。请谨慎使用。",
        parameters={
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "要删除的文件或目录路径。"},
                "scope": {"type": "string", "enum": ["workspace", "system"], "description": "路径作用域：workspace=工作空间内（默认）；system=工作空间外（需授权，filename 为绝对路径）。"},
            },
            "required": ["filename"],
        },
        permission_level=PermissionLevel.WRITE,
        handler=_delete_file,
        group="system",
        category="文件操作",
        guidance="要删除某个文件或文件夹 → delete_file",
    ))
